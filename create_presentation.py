"""Generate a Product Manager presentation for the Gamification Admin Platform."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# --- Color Palette ---
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
BLUE_700 = RGBColor(0x1D, 0x4E, 0xD8)
BLUE_500 = RGBColor(0x3B, 0x82, 0xF6)
BLUE_100 = RGBColor(0xDB, 0xEA, 0xFE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE_50 = RGBColor(0xF8, 0xFA, 0xFC)
SLATE_100 = RGBColor(0xF1, 0xF5, 0xF9)
SLATE_300 = RGBColor(0xCB, 0xD5, 0xE1)
SLATE_500 = RGBColor(0x64, 0x74, 0x8B)
SLATE_700 = RGBColor(0x33, 0x41, 0x55)
SLATE_900 = RGBColor(0x0F, 0x17, 0x2A)
EMERALD = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
VIOLET = RGBColor(0x7C, 0x3A, 0xED)
RED = RGBColor(0xDC, 0x26, 0x26)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
ORANGE = RGBColor(0xF9, 0x73, 0x16)

def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_para(tf, text, font_size=16, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, space_before=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p.alignment = alignment
    p.space_before = space_before
    return p

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, value, subtitle, accent_color):
    add_rounded_rect(slide, left, top, width, height, RGBColor(0x1E, 0x29, 0x3B))
    # accent bar
    add_rounded_rect(slide, left + 0.15, top + 0.15, 0.08, height - 0.3, accent_color)
    add_textbox(slide, left + 0.4, top + 0.2, width - 0.6, 0.3, title, 11, False, SLATE_300)
    add_textbox(slide, left + 0.4, top + 0.5, width - 0.6, 0.5, value, 28, True, WHITE)
    add_textbox(slide, left + 0.4, top + 1.0, width - 0.6, 0.3, subtitle, 10, False, SLATE_500)

def slide_number_footer(slide, num, total):
    add_textbox(slide, 11.5, 7.0, 1.5, 0.4, f"{num} / {total}", 10, False, SLATE_500, PP_ALIGN.RIGHT)

TOTAL_SLIDES = 15

# ============================================================
# SLIDE 1 — Title
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
# Accent line
add_rounded_rect(s, 1.0, 2.0, 0.8, 0.06, BLUE_500)
add_textbox(s, 1.0, 2.2, 11, 1.2, "Gamification Admin Platform", 48, True, WHITE)
add_textbox(s, 1.0, 3.5, 9, 0.8, "Loyalty + Gamification for Retail Banking", 24, False, BLUE_500)
add_textbox(s, 1.0, 4.5, 9, 0.5, "Product Overview  •  May 2026", 16, False, SLATE_300)
# decorative pill
add_rounded_rect(s, 1.0, 5.3, 2.8, 0.35, RGBColor(0x1E, 0x29, 0x3B))
add_textbox(s, 1.15, 5.32, 2.5, 0.3, "● Live Command Center", 11, True, EMERALD)
slide_number_footer(s, 1, TOTAL_SLIDES)

# ============================================================
# SLIDE 2 — Loyalty + Gamification Vision
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_textbox(s, 1.0, 0.5, 11, 0.6, "THE VISION", 12, True, BLUE_500)
add_textbox(s, 1.0, 1.0, 11, 0.8, "Loyalty Program Meets Gamification", 36, True, WHITE)

desc = ("This project combines a traditional loyalty program with gamification mechanics. "
        "Instead of passive point collection, customers actively complete quests — real banking actions — "
        "to earn XP, climb leaderboards, and qualify for a grand prize. "
        "The result: deeper engagement, stronger deposits, and measurable revenue.")
tb = add_textbox(s, 1.0, 2.0, 9, 1.5, desc, 17, False, SLATE_300)
tb.text_frame.word_wrap = True

# Three-pillar cards
points = [
    ("Loyalty", "Reward customers for\nrepeat banking behaviour\nover the full year", BLUE_500),
    ("Gamification", "Quests, XP, tiers &\nleaderboards make banking\nfun and competitive", AMBER),
    ("Grand Prize", "Year-long quest journey\nculminating in a high-value\nprize draw for top players", EMERALD),
]
for i, (title, sub, clr) in enumerate(points):
    x = 1.0 + i * 3.8
    add_rounded_rect(s, x, 4.0, 3.5, 2.2, RGBColor(0x1E, 0x29, 0x3B))
    add_rounded_rect(s, x + 0.15, 4.15, 1.5, 0.06, clr)
    add_textbox(s, x + 0.15, 4.4, 3.2, 0.4, title, 20, True, WHITE)
    add_textbox(s, x + 0.15, 4.9, 3.2, 1.2, sub, 14, False, SLATE_300)

slide_number_footer(s, 2, TOTAL_SLIDES)

# ============================================================
# SLIDE 3 — Reward Structure: Small Prizes + Grand Prize
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_textbox(s, 1.0, 0.5, 11, 0.6, "REWARD STRUCTURE", 12, True, AMBER)
add_textbox(s, 1.0, 1.0, 11, 0.8, "Redeem Small Prizes Anytime — Chase the Grand Prize All Year", 32, True, WHITE)

tb = add_textbox(s, 1.0, 1.9, 10, 0.8,
    "Customers earn XP from every quest. They can instantly redeem points for small everyday prizes "
    "OR keep accumulating XP across the full year to qualify for the grand prize draw. "
    "This dual model keeps customers engaged daily while building towards a major reward.", 16, False, SLATE_300)
tb.text_frame.word_wrap = True

# LEFT — Small Instant Prizes
add_rounded_rect(s, 1.0, 3.0, 5.3, 3.5, RGBColor(0x1E, 0x29, 0x3B))
add_rounded_rect(s, 1.15, 3.15, 2.0, 0.06, EMERALD)
add_textbox(s, 1.3, 3.35, 4.8, 0.5, "Instant Small Prizes", 20, True, WHITE)
add_textbox(s, 1.3, 3.85, 4.8, 0.35, "Redeem XP anytime for tiny rewards", 13, False, EMERALD)
small_prizes = [
    "Free coffee voucher (e.g. Starbucks)",
    "Small merchant discount coupons",
    "Mobile top-up credits",
    "Movie or event ticket discounts",
    "Cashback on next card transaction",
    "Digital gift cards (small value)",
]
for j, item in enumerate(small_prizes):
    add_textbox(s, 1.3, 4.3 + j * 0.35, 4.8, 0.3, f"●  {item}", 12, False, SLATE_300)

# RIGHT — Grand Prize
add_rounded_rect(s, 6.7, 3.0, 5.3, 3.5, RGBColor(0x1E, 0x29, 0x3B))
add_rounded_rect(s, 6.85, 3.15, 2.0, 0.06, AMBER)
add_textbox(s, 7.0, 3.35, 4.8, 0.5, "Grand Prize (Year-End)", 20, True, WHITE)
add_textbox(s, 7.0, 3.85, 4.8, 0.35, "Keep XP all year to enter the big draw", 13, False, AMBER)
grand_items = [
    "Complete many quests across the year",
    "XP accumulates toward Diamond tier",
    "Top-tier customers enter grand draw",
    "High-value prize (car, travel, cash)",
    "Creates a powerful retention loop",
    "Drives sustained banking behaviour",
]
for j, item in enumerate(grand_items):
    add_textbox(s, 7.0, 4.3 + j * 0.35, 4.8, 0.3, f"●  {item}", 12, False, SLATE_300)

# Bottom callout
add_rounded_rect(s, 1.0, 6.7, 11.3, 0.5, RGBColor(0x1E, 0x29, 0x3B))
add_textbox(s, 1.3, 6.75, 10.5, 0.4,
    "Key: Customers choose — spend points on small rewards now, or save them for the grand prize.", 14, True, BLUE_500)

slide_number_footer(s, 3, TOTAL_SLIDES)

# ============================================================
# SLIDE 4 — CASA & Savings Growth for Grand Prize
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_textbox(s, 1.0, 0.5, 11, 0.6, "DEPOSITS & SAVINGS", 12, True, EMERALD)
add_textbox(s, 1.0, 1.0, 11, 0.8, "Grow Fixed Deposits & Savings to Reach the Grand Prize", 32, True, WHITE)

tb = add_textbox(s, 1.0, 1.9, 10, 1.0,
    "A core strategy: quests encourage customers to increase their fixed deposits and savings balances. "
    "The more a customer saves, the more XP they earn — directly linking deposit growth to grand prize eligibility.", 16, False, SLATE_300)
tb.text_frame.word_wrap = True

# Two big cards side by side
for i, (title, items, clr) in enumerate([
    ("Fixed Deposit Quests", [
        "Open a new fixed deposit account",
        "Increase FD balance by 5K or more",
        "Renew an existing FD at maturity",
        "Maintain FD for 6+ months bonus XP",
    ], BLUE_500),
    ("Savings Account Quests", [
        "Save 2K+ in Al Hayrat for 2 months",
        "Set up automatic monthly savings",
        "Grow savings balance by 10% quarter",
        "Move salary to a savings account",
    ], EMERALD),
]):
    x = 1.0 + i * 5.8
    add_rounded_rect(s, x, 3.2, 5.5, 3.2, RGBColor(0x1E, 0x29, 0x3B))
    add_rounded_rect(s, x + 0.15, 3.35, 2.0, 0.06, clr)
    add_textbox(s, x + 0.25, 3.6, 5.0, 0.4, title, 20, True, WHITE)
    for j, item in enumerate(items):
        add_textbox(s, x + 0.25, 4.15 + j * 0.5, 5.0, 0.4, f"→  {item}", 14, False, SLATE_300)

add_rounded_rect(s, 1.0, 6.6, 11.3, 0.5, RGBColor(0x1E, 0x29, 0x3B))
add_textbox(s, 1.3, 6.65, 10.5, 0.4, "Result: The bank grows CASA balances while customers earn their way to the grand prize.", 14, True, AMBER)

slide_number_footer(s, 4, TOTAL_SLIDES)

# ============================================================
# SLIDE 5 — Revenue Model
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_textbox(s, 1.0, 0.5, 11, 0.6, "REVENUE MODEL", 12, True, AMBER)
add_textbox(s, 1.0, 1.0, 11, 0.8, "Multiple Revenue Streams", 36, True, WHITE)

tb = add_textbox(s, 1.0, 1.9, 10, 0.8,
    "The platform generates revenue from two major sources, creating a self-sustaining "
    "gamification ecosystem that funds prizes while driving profit.", 16, False, SLATE_300)
tb.text_frame.word_wrap = True

# Revenue source 1 — Merchants
add_rounded_rect(s, 1.0, 3.0, 5.3, 3.8, RGBColor(0x1E, 0x29, 0x3B))
add_rounded_rect(s, 1.15, 3.15, 2.0, 0.06, ORANGE)
add_textbox(s, 1.3, 3.4, 4.8, 0.5, "① Merchant Affiliates", 22, True, WHITE)
add_textbox(s, 1.3, 3.95, 4.8, 0.4, "Revenue from partner merchants", 14, False, BLUE_500)
merchant_items = [
    "Starbucks, Gulf Air, Lulu & more",
    "Commission on every card transaction",
    "Co-branded quest sponsorship fees",
    "Merchant pays for featured placement",
    "Higher volume = better merchant rates",
]
for j, item in enumerate(merchant_items):
    add_textbox(s, 1.3, 4.45 + j * 0.42, 4.8, 0.35, f"●  {item}", 13, False, SLATE_300)

# Revenue source 2 — Customer Services
add_rounded_rect(s, 6.7, 3.0, 5.3, 3.8, RGBColor(0x1E, 0x29, 0x3B))
add_rounded_rect(s, 6.85, 3.15, 2.0, 0.06, EMERALD)
add_textbox(s, 7.0, 3.4, 4.8, 0.5, "② Customer Services", 22, True, WHITE)
add_textbox(s, 7.0, 3.95, 4.8, 0.4, "Revenue when customers use bank services", 14, False, BLUE_500)
service_items = [
    "New fixed deposits & savings accounts",
    "Personal / auto / mortgage financing",
    "Bill payments & transfers via app",
    "Card activation & transaction fees",
    "Insurance & investment products",
]
for j, item in enumerate(service_items):
    add_textbox(s, 7.0, 4.45 + j * 0.42, 4.8, 0.35, f"●  {item}", 13, False, SLATE_300)

slide_number_footer(s, 5, TOTAL_SLIDES)

# ============================================================
# SLIDE 6 — Business Goals
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_textbox(s, 1.0, 0.5, 11, 0.6, "BUSINESS GOALS", 12, True, BLUE_500)
add_textbox(s, 1.0, 1.0, 11, 0.8, "What Does the Bank Achieve?", 36, True, WHITE)

goals = [
    ("Increase CASA Balances", "Encourage deposits, savings, and salary transfers", BLUE_500),
    ("Boost Digital Engagement", "Drive mobile & internet banking adoption", VIOLET),
    ("Encourage Card Spending", "Promote card usage through merchant quests", ORANGE),
    ("Reduce Risk", "Incentivize responsible banking behaviours", RED),
    ("Social Responsibility", "Support charity, sustainability & community", GREEN),
]
for i, (title, sub, clr) in enumerate(goals):
    y = 2.0 + i * 1.0
    add_rounded_rect(s, 1.0, y, 0.06, 0.7, clr)
    add_textbox(s, 1.3, y, 4.0, 0.35, title, 18, True, WHITE)
    add_textbox(s, 1.3, y + 0.35, 8.0, 0.35, sub, 13, False, SLATE_300)

slide_number_footer(s, 6, TOTAL_SLIDES)

# ============================================================
# SLIDE 4 — Target Users
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_textbox(s, 1.0, 0.5, 11, 0.6, "TARGET USERS", 12, True, BLUE_500)
add_textbox(s, 1.0, 1.0, 11, 0.8, "Who Uses the Platform?", 36, True, WHITE)

users = [
    "Campaign Managers — create & manage quests",
    "Product Owners — define objectives & segments",
    "Digital Banking Team — publish to channels",
    "Marketing Team — plan promotions & campaigns",
    "Risk & Compliance — review rules & behaviour",
    "Management — monitor KPIs, ROI & performance",
]
for i, u in enumerate(users):
    y = 2.2 + i * 0.75
    add_rounded_rect(s, 1.0, y, 10.5, 0.6, RGBColor(0x1E, 0x29, 0x3B))
    add_textbox(s, 1.3, y + 0.1, 9.5, 0.4, u, 16, False, WHITE)

slide_number_footer(s, 7, TOTAL_SLIDES)

# ============================================================
# SLIDE 5 — Quest Categories (5 categories)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_textbox(s, 1.0, 0.5, 11, 0.6, "STRATEGIC CATEGORIES", 12, True, BLUE_500)
add_textbox(s, 1.0, 1.0, 11, 0.8, "Five Pillars of Gamification", 36, True, WHITE)

cats = [
    ("Improve CASA", "Salary movement, savings,\ndeposit growth", "e.g. Salary Anchor", BLUE_500),
    ("Increase Engagement", "App usage, bill pay,\ndigital adoption", "e.g. Bill Pay Streak", VIOLET),
    ("Encourage Spending", "Card transactions,\nmerchant partnerships", "e.g. Starbucks Quest", ORANGE),
    ("Risk Reduction", "On-time payments,\nphishing awareness", "e.g. Phishing Fighter", RED),
    ("Social Responsibility", "Charity, green finance,\ncommunity impact", "e.g. Donate to Charity", GREEN),
]
for i, (title, desc, ex, clr) in enumerate(cats):
    x = 0.5 + i * 2.5
    add_rounded_rect(s, x, 2.2, 2.3, 3.8, RGBColor(0x1E, 0x29, 0x3B))
    add_rounded_rect(s, x + 0.15, 2.35, 1.9, 0.06, clr)
    add_textbox(s, x + 0.2, 2.6, 1.9, 0.5, title, 15, True, WHITE)
    add_textbox(s, x + 0.2, 3.2, 1.9, 1.0, desc, 12, False, SLATE_300)
    add_textbox(s, x + 0.2, 4.8, 1.9, 0.5, ex, 11, True, clr)

slide_number_footer(s, 8, TOTAL_SLIDES)

# ============================================================
# SLIDE 6 — User Journey
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_textbox(s, 1.0, 0.5, 11, 0.6, "USER JOURNEY", 12, True, BLUE_500)
add_textbox(s, 1.0, 1.0, 11, 0.8, "End-to-End Workflow", 36, True, WHITE)

steps = [
    "1. Create a new quest via Quest Forge form",
    "2. Set category, dates, segment & financials",
    "3. Quest appears in Campaign Calendar",
    "4. Quest listed in Digital Channel Table",
    "5. Monitor customer participation & rankings",
    "6. Track KPIs, revenue forecast & ROI",
    "7. Export data as CSV for reporting",
]
for i, step in enumerate(steps):
    y = 2.2 + i * 0.7
    add_rounded_rect(s, 1.0, y, 0.06, 0.5, BLUE_500)
    add_textbox(s, 1.3, y + 0.05, 10, 0.4, step, 17, False, WHITE)

slide_number_footer(s, 9, TOTAL_SLIDES)

# ============================================================
# SLIDE 7 — Dashboard KPIs
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_textbox(s, 1.0, 0.5, 11, 0.6, "COMMAND CENTER", 12, True, BLUE_500)
add_textbox(s, 1.0, 1.0, 11, 0.8, "Dashboard KPIs at a Glance", 36, True, WHITE)

add_card(s, 1.0, 2.2, 2.7, 1.4, "Active / Scheduled", "13", "total quest programs", BLUE_500)
add_card(s, 4.0, 2.2, 2.7, 1.4, "Customers Doing Quests", "249K", "forecast completion", EMERALD)
add_card(s, 7.0, 2.2, 2.7, 1.4, "Forecast Revenue", "$1.7M", "across all quests", AMBER)
add_card(s, 10.0, 2.2, 2.7, 1.4, "Forecast ROI", "594%", "after prize budget", VIOLET)

add_textbox(s, 1.0, 4.2, 11, 0.5, "These KPIs auto-calculate from quest data — no manual input needed.", 14, False, SLATE_500)

# Category summary
add_textbox(s, 1.0, 5.0, 11, 0.5, "Category Summary Cards", 20, True, WHITE)
add_textbox(s, 1.0, 5.5, 11, 0.8, "Each of the 5 categories shows: number of quests, participating customers, and ROI.\nThis helps PMs ensure a balanced portfolio across all strategic objectives.", 14, False, SLATE_300)

slide_number_footer(s, 10, TOTAL_SLIDES)

# ============================================================
# SLIDE 8 — Key Features Overview
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_textbox(s, 1.0, 0.5, 11, 0.6, "KEY FEATURES", 12, True, BLUE_500)
add_textbox(s, 1.0, 1.0, 11, 0.8, "Platform Capabilities", 36, True, WHITE)

features = [
    ("Quest Forge Form", "Create/edit quests with title, category, dates,\nsegment, targets, XP, revenue & budget", BLUE_500),
    ("Campaign Calendar", "Month/week views, filters by category & status,\nload indicators, quick-schedule with + button", VIOLET),
    ("Digital Channel Table", "Full quest list with search, sort, filter, paginate,\nedit, delete, and CSV export", EMERALD),
    ("Customer Ranking", "Top 10 leaderboard by XP, completion & impact;\nDiamond / Platinum / Gold / Silver tiers", AMBER),
    ("Analytics Suite", "Revenue forecast, quest mix donut, customer\nfunnel, and strategic KPI radar scorecard", ORANGE),
    ("Data Persistence", "Quests stored in browser localStorage;\nreload sample data with one click", SLATE_300),
]
for i, (title, desc, clr) in enumerate(features):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.0
    y = 2.2 + row * 2.4
    add_rounded_rect(s, x, y, 3.7, 2.1, RGBColor(0x1E, 0x29, 0x3B))
    add_rounded_rect(s, x + 0.15, y + 0.15, 0.08, 1.8, clr)
    add_textbox(s, x + 0.4, y + 0.2, 3.1, 0.4, title, 17, True, WHITE)
    add_textbox(s, x + 0.4, y + 0.7, 3.1, 1.2, desc, 13, False, SLATE_300)

slide_number_footer(s, 11, TOTAL_SLIDES)

# ============================================================
# SLIDE 9 — Sample Quests
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_textbox(s, 1.0, 0.5, 11, 0.6, "SAMPLE QUESTS", 12, True, BLUE_500)
add_textbox(s, 1.0, 1.0, 11, 0.8, "Pre-loaded Campaign Examples", 36, True, WHITE)

sample_quests = [
    ("Salary Anchor", "CASA", "Move salary to the bank", "42K customers"),
    ("Bill Pay Streak", "ENGAGE", "Pay 2 bills via mobile", "55K customers"),
    ("Navigate in the App", "SPEND", "Complete guided app tutorial", "36K customers"),
    ("Starbucks Coffee", "SPEND", "Buy coffee with bank card", "28K customers"),
    ("Phishing Fighter", "RISK", "Report phishing attempts", "75K customers"),
    ("Pay CC On Time", "RISK", "Timely credit card payment", "75K customers"),
    ("Donate to Charity", "SOCIAL", "Donate to verified charity", "22K customers"),
    ("Green Financing", "SOCIAL", "Finance a solar project", "18K customers"),
]
# Header
add_rounded_rect(s, 1.0, 2.0, 11.3, 0.45, BLUE_700)
add_textbox(s, 1.2, 2.05, 3.0, 0.35, "Quest", 12, True, WHITE)
add_textbox(s, 4.2, 2.05, 1.5, 0.35, "Category", 12, True, WHITE)
add_textbox(s, 5.8, 2.05, 4.0, 0.35, "Customer Action", 12, True, WHITE)
add_textbox(s, 10.0, 2.05, 2.0, 0.35, "Target", 12, True, WHITE)

for i, (name, cat, action, target) in enumerate(sample_quests):
    y = 2.5 + i * 0.55
    bg_clr = RGBColor(0x1E, 0x29, 0x3B) if i % 2 == 0 else RGBColor(0x17, 0x22, 0x34)
    add_rounded_rect(s, 1.0, y, 11.3, 0.5, bg_clr)
    add_textbox(s, 1.2, y + 0.07, 3.0, 0.35, name, 13, True, WHITE)
    add_textbox(s, 4.2, y + 0.07, 1.5, 0.35, cat, 12, False, BLUE_500)
    add_textbox(s, 5.8, y + 0.07, 4.0, 0.35, action, 12, False, SLATE_300)
    add_textbox(s, 10.0, y + 0.07, 2.0, 0.35, target, 12, False, EMERALD)

slide_number_footer(s, 12, TOTAL_SLIDES)

# ============================================================
# SLIDE 10 — Operating Model
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_textbox(s, 1.0, 0.5, 11, 0.6, "OPERATING MODEL", 12, True, BLUE_500)
add_textbox(s, 1.0, 1.0, 11, 0.8, "Suggested Roles & Responsibilities", 36, True, WHITE)

roles = [
    ("Campaign Manager", "Creates and manages quests day-to-day"),
    ("Product Owner", "Defines business objectives and target segments"),
    ("Risk Team", "Reviews risk quests and behaviour rules"),
    ("Compliance Team", "Reviews prize-linked campaign regulations"),
    ("Digital Team", "Publishes quests to mobile and web channels"),
    ("Management", "Reviews KPIs, ROI, and program performance"),
]
for i, (role, resp) in enumerate(roles):
    y = 2.2 + i * 0.85
    add_rounded_rect(s, 1.0, y, 11, 0.7, RGBColor(0x1E, 0x29, 0x3B))
    add_textbox(s, 1.3, y + 0.1, 3.0, 0.5, role, 17, True, BLUE_500)
    add_textbox(s, 4.5, y + 0.15, 7.0, 0.4, resp, 15, False, SLATE_300)

slide_number_footer(s, 13, TOTAL_SLIDES)

# ============================================================
# SLIDE 11 — Roadmap to Production
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_textbox(s, 1.0, 0.5, 11, 0.6, "ROADMAP TO PRODUCTION", 12, True, BLUE_500)
add_textbox(s, 1.0, 1.0, 11, 0.8, "What's Needed Before Go-Live", 36, True, WHITE)

add_textbox(s, 1.0, 2.0, 5.5, 0.5, "Current State: Front-end prototype (HTML + JS)", 14, False, AMBER)

items_left = [
    "User authentication & access control",
    "Secure database storage",
    "Approval workflow (maker-checker)",
    "Audit trail for edits & deletions",
    "Role-based permissions",
]
items_right = [
    "Mobile & internet banking integration",
    "Customer data & transaction APIs",
    "Compliance review for prize rules",
    "Data privacy controls (GDPR/PDPL)",
    "Performance & load testing",
]

add_textbox(s, 1.0, 2.8, 5, 0.4, "Security & Governance", 18, True, WHITE)
for i, item in enumerate(items_left):
    add_rounded_rect(s, 1.0, 3.3 + i * 0.65, 0.06, 0.45, RED)
    add_textbox(s, 1.3, 3.35 + i * 0.65, 5, 0.35, item, 14, False, SLATE_300)

add_textbox(s, 7.0, 2.8, 5, 0.4, "Integration & Compliance", 18, True, WHITE)
for i, item in enumerate(items_right):
    add_rounded_rect(s, 7.0, 3.3 + i * 0.65, 0.06, 0.45, BLUE_500)
    add_textbox(s, 7.3, 3.35 + i * 0.65, 5, 0.35, item, 14, False, SLATE_300)

slide_number_footer(s, 14, TOTAL_SLIDES)

# ============================================================
# SLIDE 12 — Summary & Next Steps
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_rounded_rect(s, 1.0, 2.0, 0.8, 0.06, BLUE_500)
add_textbox(s, 1.0, 2.2, 11, 1.0, "Summary", 42, True, WHITE)

summary_text = ("This platform turns customer banking actions into structured campaigns with points, "
                "rewards, monitoring, and management reporting. It is a single dashboard for creating, "
                "scheduling, tracking, and analysing prize-linked quests across all digital channels.")
tb = add_textbox(s, 1.0, 3.3, 9, 1.2, summary_text, 17, False, SLATE_300)
tb.text_frame.word_wrap = True

add_textbox(s, 1.0, 4.8, 11, 0.5, "Next Steps", 22, True, BLUE_500)
next_steps = [
    "✓  Review prototype with stakeholders",
    "✓  Prioritize production requirements",
    "✓  Define integration scope with IT",
    "✓  Plan pilot with one quest category",
]
for i, step in enumerate(next_steps):
    add_textbox(s, 1.0, 5.4 + i * 0.45, 10, 0.4, step, 16, False, WHITE)

add_textbox(s, 1.0, 7.0, 11, 0.4, "Thank you", 14, True, SLATE_500)
slide_number_footer(s, 15, TOTAL_SLIDES)

# ============================================================
# Save
# ============================================================
output_path = os.path.join(os.path.dirname(__file__), "Gamification_Admin_Platform_Overview.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
